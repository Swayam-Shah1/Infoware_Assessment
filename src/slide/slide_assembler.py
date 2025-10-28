"""
Slide Assembler
Creates PowerPoint presentations from slide content
"""
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from utils.config_manager import ConfigManager
from utils.logger import Logger
import os


class SlideAssembler:
    """Assembles PowerPoint presentations"""
    
    def __init__(self, config: ConfigManager, logger: Logger):
        """
        Initialize slide assembler
        
        Args:
            config: Configuration manager
            logger: Logger instance
        """
        self.config = config
        self.logger = logger
        
        self.aspect_ratio = config.get_value('slides.aspect_ratio', '16:9')
        self.theme = config.get_value('slides.theme', 'modern')
    
    def create_presentation(self, slides: List[Dict[str, Any]], output_path: str) -> None:
        """
        Create PowerPoint presentation from slides
        
        Args:
            slides: List of enriched SlideContent dictionaries
            output_path: Output file path
        """
        self.logger.info(f"Creating presentation with {len(slides)} slides...")
        
        # Create presentation
        prs = Presentation()
        
        # Set slide size based on aspect ratio
        if self.aspect_ratio == '16:9':
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
        else:  # 4:3
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        
        # Add slides
        for slide_data in slides:
            self._add_slide(prs, slide_data)
        
        # Save presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        
        self.logger.info(f"Presentation saved to: {output_path}")
    
    def _add_slide(self, prs: Presentation, slide_data: Dict[str, Any]) -> None:
        """
        Add a single slide to presentation
        
        Args:
            prs: PowerPoint presentation
            slide_data: SlideContent dictionary
        """
        # Create blank slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xF8)  # Light blue
        
        # Add title with proper sizing and word wrapping
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        # Set title text
        title_text = slide_data.get('title', 'Untitled')
        title_frame.text = title_text
        
        # Format title paragraph with dynamic font sizing
        title_para = title_frame.paragraphs[0]
        
        # Calculate dynamic font size based on title length
        title_length = len(title_text)
        word_count = len(title_text.split())
        
        if title_length > 60 or word_count > 8:
            # Long titles: smaller font
            title_font_size = Pt(24)
        elif title_length > 40 or word_count > 6:
            # Medium titles
            title_font_size = Pt(28)
        else:
            # Short titles: larger font
            title_font_size = Pt(32)
        
        title_para.font.size = title_font_size
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x4A, 0x90, 0xE2)  # Blue color
        title_para.alignment = PP_ALIGN.LEFT
        title_para.space_before = Pt(0)
        title_para.space_after = Pt(12)
        
        # Enable auto-fit to prevent text overflow
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Check if there's an image to determine layout
        has_image = 'visual' in slide_data and 'path' in slide_data['visual'] and os.path.exists(slide_data['visual']['path'])
        
        # Add bullets with proper formatting (positioned below title)
        # Adjust width based on whether there's an image
        if has_image:
            # Two-column layout: content on left, image on right
            bullet_box_width = Inches(4.5)
            image_x = Inches(5.5)
            left_margin = Inches(0.5)
        else:
            # Full-width layout: content spans entire width
            bullet_box_width = Inches(9)
            left_margin = Inches(0.5)
        
        # Position content box below title (start at 1.5 inches from top)
        bullet_box = slide.shapes.add_textbox(left_margin, Inches(1.5), bullet_box_width, Inches(3.8))
        bullet_frame = bullet_box.text_frame
        bullet_frame.word_wrap = True
        bullet_frame.margin_left = Inches(0.1)   # Internal left margin
        bullet_frame.margin_right = Inches(0.1)  # Internal right margin
        bullet_frame.margin_top = Inches(0.05)   # Internal top margin
        bullet_frame.margin_bottom = Inches(0.05)  # Internal bottom margin
        
        bullets = slide_data.get('bullets', [])
        
        # Clear existing paragraphs
        while len(bullet_frame.paragraphs) > 0:
            bullet_frame._element.remove(bullet_frame.paragraphs[0]._element)
        
        # Add bullets as separate paragraphs with proper spacing
        for i, bullet_text in enumerate(bullets):
            # Clean up bullet text (remove extra characters and gibberish)
            clean_bullet = bullet_text.strip()
            
            # Skip empty or very short bullets (likely gibberish)
            if len(clean_bullet) < 10:
                continue
            
            # Create new paragraph for each bullet
            para = bullet_frame.add_paragraph()
            
            # Format the paragraph first
            para.alignment = PP_ALIGN.LEFT
            para.level = 0
            para.space_before = Pt(0)
            para.space_after = Pt(10)  # Space between bullets
            
            # Use simple bullet point formatting (no complex indents)
            para.text = f"• {clean_bullet}"
            
            # Calculate dynamic font size based on bullet length
            bullet_length = len(clean_bullet)
            word_count = len(clean_bullet.split())
            
            # Determine font size based on content length
            if bullet_length > 200 or word_count > 20:
                bullet_font_size = Pt(11)  # Very long bullets
            elif bullet_length > 150 or word_count > 15:
                bullet_font_size = Pt(12)  # Long bullets
            elif bullet_length > 100 or word_count > 12:
                bullet_font_size = Pt(13)  # Medium bullets
            else:
                bullet_font_size = Pt(14)  # Short bullets
            
            # Format the entire paragraph
            para.font.size = bullet_font_size
            para.font.color.rgb = RGBColor(0x22, 0x22, 0x22)  # Dark gray
            para.font.name = 'Calibri'
            
            # Use simple indentation - just set level for bullets
            para.level = 0  # Use level 0 to avoid extra indentation
        
        # Add image if available (on right side if exists)
        if has_image:
            img_path = slide_data['visual']['path']
            try:
                # Add image to right side, aligned with content
                slide.shapes.add_picture(
                    img_path,
                    image_x, Inches(1.5),  # Match content Y position
                    width=Inches(4), height=Inches(3.8)  # Match content height
                )
            except Exception as e:
                self.logger.warning(f"Could not add image {img_path}: {e}")

