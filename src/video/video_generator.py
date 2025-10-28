"""
Video Generator
Creates video from PowerPoint presentation
"""
import os
import tempfile
from typing import List, Dict, Any, Optional
from utils.config_manager import ConfigManager
from utils.logger import Logger


class VideoGenerator:
    """Generates video from presentation"""
    
    def __init__(self, config: ConfigManager, logger: Logger):
        """
        Initialize video generator
        
        Args:
            config: Configuration manager
            logger: Logger instance
        """
        self.config = config
        self.logger = logger
        
        # Get video settings
        self.output_format = config.get_value('video.output_format', 'mp4')
        self.fps = config.get_value('video.fps', 30)
        self.resolution_width = config.get_value('video.resolution.width', 1920)
        self.resolution_height = config.get_value('video.resolution.height', 1080)
        self.min_duration = config.get_value('video.min_slide_duration', 5)
        self.max_duration = config.get_value('video.max_slide_duration', 12)
        self.ken_burns = config.get_value('video.ken_burns_enabled', True)
        self.tts_rate = config.get_value('video.tts_rate', 150)
    
    def create_video(self, presentation_path: str, output_path: str) -> None:
        """
        Create video from presentation
        
        Args:
            presentation_path: Path to PPTX file
            output_path: Output video path
        """
        self.logger.info(f"Generating video from: {presentation_path}")
        self.logger.info(f"Output format requested: {self.output_format}")
        
        try:
            # Extract slide content from presentation
            slide_content = self._extract_slide_content(presentation_path)
            
            # Create video from actual slide content
            success = self._create_video_from_slides(slide_content, output_path)
            
            if not success:
                raise Exception("Video creation failed")
                
            self.logger.info(f"Video saved to: {output_path}")
        
        except Exception as e:
            self.logger.error(f"Video generation failed: {e}")
            raise
    
    def _extract_slide_content(self, presentation_path: str) -> List[Dict[str, Any]]:
        """
        Extract slide content from PowerPoint presentation
        
        Args:
            presentation_path: Path to PPTX file
            
        Returns:
            List of slide content dictionaries
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(presentation_path)
            slides_data = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    'index': slide_idx,
                    'title': '',
                    'content': []
                }
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if not slide_data['title']:
                            # First text shape is usually the title
                            slide_data['title'] = text
                        else:
                            # Other text shapes are content
                            if text:
                                slide_data['content'].append(text)
                
                slides_data.append(slide_data)
                self.logger.debug(f"Extracted slide {slide_idx}: Title='{slide_data['title']}', Content lines={len(slide_data['content'])}")
            
            self.logger.info(f"Extracted content from {len(slides_data)} slides")
            return slides_data
            
        except Exception as e:
            self.logger.error(f"Failed to extract slide content: {e}")
            return []
    
    def _create_video_from_slides(self, slide_content: List[Dict[str, Any]], output_path: str) -> bool:
        """
        Create video from extracted slide content with TTS
        
        Args:
            slide_content: List of slide content dictionaries
            output_path: Output video path
            
        Returns:
            True if successful
        """
        if not slide_content:
            self.logger.error("No slide content available to create video")
            return False
        
        try:
            from PIL import Image, ImageDraw, ImageFont
            import tempfile
            
            frames_dir = tempfile.mkdtemp(prefix='video_frames_')
            audio_files = []
            
            # Generate TTS audio for each slide
            tts_audio_generated = False
            for slide_idx, slide_data in enumerate(slide_content):
                title = slide_data.get('title', '')
                content_text = '\n'.join(slide_data.get('content', []))
                
                # Use speaker notes if available, otherwise use bullets
                speaker_notes = slide_data.get('speaker_notes', '')
                
                # Create full text for narration
                if speaker_notes:
                    # Use speaker notes for more natural narration
                    full_text = f"{title}. {speaker_notes}"
                else:
                    # Fallback to title + bullets
                    full_text = title
                    if content_text:
                        full_text += f". {content_text}"
                
                if full_text:
                    audio_path = self._generate_tts_audio(full_text, slide_idx, frames_dir)
                    if audio_path:
                        audio_files.append(audio_path)
                        tts_audio_generated = True
            
            if tts_audio_generated:
                self.logger.info(f"TTS audio generated for {len(audio_files)} slides")
            else:
                self.logger.warning("TTS not available")
            
            # Calculate actual audio durations to sync with video
            audio_durations = []
            for slide_idx, slide_data in enumerate(slide_content):
                if slide_idx < len(audio_files) and audio_files[slide_idx] and os.path.exists(audio_files[slide_idx]):
                    # Get audio duration from the file
                    import wave
                    try:
                        with wave.open(audio_files[slide_idx], 'rb') as wf:
                            audio_duration = wf.getnframes() / wf.getframerate()
                            audio_durations.append(max(audio_duration, self.min_duration))
                    except Exception as e:
                        self.logger.warning(f"Could not get audio duration for slide {slide_idx}: {e}")
                        audio_durations.append(self.min_duration)
                else:
                    audio_durations.append(self.min_duration)
            
            # Create frames for each slide using actual audio durations
            frames_list = []
            for slide_idx, slide_data in enumerate(slide_content):
                title = slide_data.get('title', 'Slide ' + str(slide_idx + 1))
                content_lines = slide_data.get('content', [])
                
                # Use actual audio duration for this slide to ensure sync
                slide_duration = audio_durations[slide_idx] if slide_idx < len(audio_durations) else self.min_duration
                
                # Generate frames WITHOUT captions initially
                slide_frames = self._create_slide_frames(
                    title, content_lines, slide_duration, self.fps, show_captions=False
                )
                
                # Save frames to disk
                for frame_idx, frame_img in enumerate(slide_frames):
                    frame_path = os.path.join(frames_dir, f"slide_{slide_idx:02d}_frame_{frame_idx:04d}.png")
                    frame_img.save(frame_path)
                    frames_list.append(frame_path)
            
            # Create video using imageio which handles frame rates better
            success = self._create_video_with_imageio(frames_list, output_path, audio_files)
            
            # If video creation failed, recreate frames WITH captions
            if not success and tts_audio_generated:
                self.logger.warning("Audio merge failed - recreating video with on-screen captions")
                
                # Clear old frames
                for frame_path in frames_list:
                    if os.path.exists(frame_path):
                        os.remove(frame_path)
                frames_list = []
                
                # Recreate frames with captions (use same audio durations for consistency)
                for slide_idx, slide_data in enumerate(slide_content):
                    title = slide_data.get('title', 'Slide ' + str(slide_idx + 1))
                    content_lines = slide_data.get('content', [])
                    
                    # Use the same audio duration as before
                    slide_duration = audio_durations[slide_idx] if slide_idx < len(audio_durations) else self.min_duration
                    
                    # Generate frames WITH captions
                    slide_frames = self._create_slide_frames(
                        title, content_lines, slide_duration, self.fps, show_captions=True
                    )
                    
                    # Save frames to disk
                    for frame_idx, frame_img in enumerate(slide_frames):
                        frame_path = os.path.join(frames_dir, f"slide_{slide_idx:02d}_frame_{frame_idx:04d}.png")
                        frame_img.save(frame_path)
                        frames_list.append(frame_path)
                
                # Create video again with captions (without audio)
                success = self._create_video_with_imageio(frames_list, output_path, [])
            
            # Cleanup
            import shutil
            if os.path.exists(frames_dir):
                shutil.rmtree(frames_dir)
            
            return success
            
        except Exception as e:
            self.logger.error(f"Video creation from slides failed: {e}")
            import traceback
            self.logger.debug(traceback.format_exc())
            return False
    
    def _generate_tts_audio(self, text: str, slide_idx: int, output_dir: str) -> Optional[str]:
        """
        Generate TTS audio for text
        
        Args:
            text: Text to convert to speech
            slide_idx: Slide index
            output_dir: Output directory for audio
            
        Returns:
            Path to audio file or None
        """
        try:
            import pyttsx3
            
            audio_path = os.path.join(output_dir, f"audio_{slide_idx:02d}.wav")
            
            try:
                engine = pyttsx3.init()
                
                # Configure voice
                engine.setProperty('rate', self.tts_rate)
                engine.setProperty('volume', 0.8)
                
                # Generate audio
                engine.save_to_file(text, audio_path)
                engine.runAndWait()
                
                if os.path.exists(audio_path) and os.path.getsize(audio_path) > 0:
                    self.logger.debug(f"Generated TTS audio for slide {slide_idx}: {len(text)} chars")
                    return audio_path
                else:
                    self.logger.warning(f"TTS audio generation produced empty file for slide {slide_idx}")
                    return None
                    
            except Exception as e:
                self.logger.warning(f"TTS generation failed for slide {slide_idx}: {e}")
                return None
                
        except Exception as e:
            self.logger.warning(f"TTS not available: {e}")
            return None
    
    def _estimate_video_duration(self, audio_files: List[str], num_slides: int) -> float:
        """
        Estimate total video duration
        
        Args:
            audio_files: List of audio file paths
            num_slides: Number of slides
            
        Returns:
            Estimated duration in seconds
        """
        if audio_files:
            # Try to get duration from audio files
            try:
                total_duration = 0
                for audio_file in audio_files:
                    if os.path.exists(audio_file):
                        # Estimate ~200 words per minute (average speech rate)
                        # Approximate duration based on file size or use default
                        total_duration += self.min_duration
            except:
                pass
        
        # Default duration: 5-8 seconds per slide
        return num_slides * self.min_duration
    
    def _estimate_slide_duration(self, content_lines: List[str]) -> float:
        """
        Estimate duration for a slide based on content
        
        Args:
            content_lines: List of content lines
            
        Returns:
            Duration in seconds
        """
        total_chars = sum(len(line) for line in content_lines)
        
        # Estimate based on text length: ~150 words per minute
        estimated_words = total_chars / 5
        estimated_seconds = (estimated_words / 150) * 60
        
        # Clamp between min and max duration
        return max(self.min_duration, min(self.max_duration, estimated_seconds))
    
    def _create_slide_frames(self, title: str, content_lines: List[str], duration: float, fps: int, show_captions: bool = True) -> List:
        """
        Create visual frames for a slide
        
        Args:
            title: Slide title
            content_lines: List of content lines
            duration: Slide duration in seconds
            fps: Frames per second
            show_captions: Whether to show text captions on frames
            
        Returns:
            List of PIL Image objects
        """
        from PIL import Image, ImageDraw, ImageFont
        
        frames = []
        num_frames = int(duration * fps)
        
        # Load font
        try:
            title_font = ImageFont.truetype("arial.ttf", 64)
            content_font = ImageFont.truetype("arial.ttf", 32)
            small_font = ImageFont.truetype("arial.ttf", 28)
        except:
            title_font = ImageFont.load_default()
            content_font = ImageFont.load_default()
            small_font = ImageFont.load_default()
        
        for i in range(num_frames):
            # Create frame
            frame = Image.new('RGB', (self.resolution_width, self.resolution_height), color='#E8F4F8')
            draw = ImageDraw.Draw(frame)
            
            # Draw title
            bbox = draw.textbbox((0, 0), title, font=title_font)
            title_width = bbox[2] - bbox[0]
            title_x = (self.resolution_width - title_width) // 2
            draw.text((title_x, 150), title, fill='#4A90E2', font=title_font)
            
            # Draw content lines - properly formatted with bullets and line breaks
            y_offset = 350
            line_height = 70  # Increased spacing between items
            bullet_height = 60  # Spacing within a bullet item
            max_width = self.resolution_width - 600  # Left and right margins
            x_start = 250  # Starting position
            
            # Space between bullet items
            item_spacing = 30
            
            for idx, line in enumerate(content_lines):
                if not line.strip() or y_offset >= self.resolution_height - 100:
                    continue
                
                # Check if this is a bullet point (starts with • or - or is a bullet)
                is_bullet = line.strip().startswith('•') or line.strip().startswith('-') or line.strip().startswith('\u2022')
                
                # Extract bullet symbol if present
                bullet_symbol = ''
                if is_bullet:
                    clean_line = line.strip()[1:].strip()  # Remove first character
                    bullet_symbol = '• '
                else:
                    clean_line = line.strip()
                
                # Add spacing between items
                if idx > 0:
                    y_offset += item_spacing
                
                if clean_line:
                    # Word wrap if line is too long
                    words = clean_line.split()
                    wrapped_lines = []
                    current_line = []
                    current_width = 0
                    
                    for word in words:
                        word_bbox = draw.textbbox((0, 0), word, font=content_font)
                        word_width = word_bbox[2] - word_bbox[0]
                        space_width = draw.textbbox((0, 0), ' ', font=content_font)[2] - draw.textbbox((0, 0), ' ', font=content_font)[0]
                        
                        if current_width + word_width + space_width < max_width:
                            current_line.append(word)
                            current_width += word_width + space_width
                        else:
                            if current_line:
                                wrapped_lines.append(' '.join(current_line))
                            current_line = [word]
                            current_width = word_width
                    
                    if current_line:
                        wrapped_lines.append(' '.join(current_line))
                    
                    # Draw wrapped lines with bullet
                    for line_idx, wrapped_line in enumerate(wrapped_lines):
                        if y_offset < self.resolution_height - 100:
                            # For first line of bullet, add bullet symbol
                            if line_idx == 0 and bullet_symbol:
                                # Draw bullet symbol
                                bullet_bbox = draw.textbbox((0, 0), bullet_symbol, font=content_font)
                                draw.text((x_start, y_offset), bullet_symbol, fill='#4A90E2', font=content_font)
                                # Draw text starting after bullet
                                text_x = x_start + (bullet_bbox[2] - bullet_bbox[0]) + 10
                                draw.text((text_x, y_offset), wrapped_line, fill='#2C3E50', font=content_font)
                            else:
                                # Draw wrapped line without bullet (indented for continuation)
                                draw.text((x_start + 30, y_offset), wrapped_line, fill='#2C3E50', font=content_font)
                            
                            # Use different spacing for first line vs wrapped lines
                            y_offset += bullet_height if line_idx == 0 else (bullet_height - 15)
            
            # Add caption indicator at bottom
            if show_captions:
                # Show text on all frames when no TTS
                if i < num_frames // 4:  # Show in first quarter of frames
                    caption_text = "Text to Speech unavailable - displaying text captions"
                else:
                    caption_text = "View the slide content above"
                
                bbox = draw.textbbox((0, 0), caption_text, font=small_font)
                caption_width = bbox[2] - bbox[0]
                caption_x = (self.resolution_width - caption_width) // 2
                caption_y = self.resolution_height - 60
                
                # Draw background for caption
                padding = 10
                draw.rectangle(
                    [caption_x - padding, caption_y - padding, 
                     caption_x + caption_width + padding, caption_y + 40 + padding],
                    fill='#FFFFFF',
                    outline='#4A90E2',
                    width=2
                )
                
                draw.text((caption_x, caption_y), caption_text, fill='#2C3E50', font=small_font)
            
            frames.append(frame)
        
        return frames
    
    
    def _create_video_with_imageio(self, frames_list: List[str], output_path: str, audio_files: List[str]) -> bool:
        """
        Create video with imageio and add audio with ffmpeg
        
        Args:
            frames_list: List of frame paths
            output_path: Output video path
            audio_files: List of audio paths
            
        Returns:
            True if successful
        """
        try:
            import imageio
            import subprocess
            import tempfile
            
            # Create video from frames using imageio
            self.logger.info(f"Creating video with imageio from {len(frames_list)} frames...")
            
            # Read all images into memory
            images = []
            for frame_path in frames_list:
                if os.path.exists(frame_path):
                    img = imageio.imread(frame_path)
                    images.append(img)
            
            # Create temporary video file
            temp_video = output_path.replace('.mp4', '_temp_video.mp4')
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # Write video using imageio
            imageio.mimwrite(temp_video, images, fps=self.fps, codec='libx264')
            
            if not os.path.exists(temp_video) or os.path.getsize(temp_video) < 1000:
                self.logger.error("Failed to create temp video with imageio")
                return False
            
            # If audio files exist, merge them
            if audio_files and len(audio_files) > 0:
                return self._merge_audio_to_video(temp_video, audio_files, output_path)
            else:
                # No audio, just move temp video to output
                import shutil
                shutil.move(temp_video, output_path)
                return True
                
        except Exception as e:
            self.logger.error(f"Failed to create video with imageio: {e}")
            import traceback
            self.logger.debug(traceback.format_exc())
            return False
    
    def _merge_audio_to_video(self, temp_video: str, audio_files: List[str], output_path: str) -> bool:
        """
        Merge audio files to video
        
        Args:
            temp_video: Path to video without audio
            audio_files: List of audio file paths
            output_path: Final output path
            
        Returns:
            True if successful
        """
        try:
            import subprocess
            import tempfile
            import imageio_ffmpeg
            
            ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
            
            # Concatenate audio files
            temp_audio = output_path.replace('.mp4', '_merged_audio.wav')
            audio_list_file = tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8')
            
            for audio_file in audio_files:
                if os.path.exists(audio_file):
                    abs_path = os.path.abspath(audio_file).replace('\\', '/')
                    audio_list_file.write(f"file '{abs_path}'\n")
            
            audio_list_file.close()
            
            # Concatenate
            concat_cmd = [
                ffmpeg_exe, '-y',
                '-f', 'concat',
                '-safe', '0',
                '-i', audio_list_file.name,
                '-acodec', 'pcm_s16le',
                '-ar', '22050',
                temp_audio
            ]
            
            self.logger.info(f"Concatenating {len(audio_files)} audio files to: {temp_audio}")
            
            concat_result = subprocess.run(concat_cmd, capture_output=True, text=True, timeout=60)
            
            if concat_result.returncode != 0:
                self.logger.error(f"Audio concatenation failed: {concat_result.stderr[:500]}")
                if os.path.exists(audio_list_file.name):
                    os.remove(audio_list_file.name)
                return False
                
            if not os.path.exists(temp_audio):
                self.logger.error("Concatenated audio file not created")
                if os.path.exists(audio_list_file.name):
                    os.remove(audio_list_file.name)
                return False
                
            self.logger.info(f"Audio concatenated successfully: {os.path.getsize(temp_audio)} bytes")
            
            # Merge audio with video
            # Don't use -shortest, let video determine the duration
            # Use map to ensure both streams are included
            merge_cmd = [
                ffmpeg_exe, '-y',
                '-i', temp_video,
                '-i', temp_audio,
                '-c:v', 'copy',
                '-c:a', 'aac',
                '-map', '0:v:0',
                '-map', '1:a:0',
                output_path
            ]
            
            merge_result = subprocess.run(merge_cmd, capture_output=True, text=True, timeout=120)
            
            # Log merge result
            if merge_result.returncode != 0:
                self.logger.error(f"Audio merge failed: {merge_result.stderr[:500]}")
            else:
                self.logger.info(f"Audio merge completed successfully")
            
            # Cleanup
            if os.path.exists(audio_list_file.name):
                os.remove(audio_list_file.name)
            if os.path.exists(temp_audio):
                os.remove(temp_audio)
            if os.path.exists(temp_video):
                os.remove(temp_video)
            
            if merge_result.returncode == 0 and os.path.exists(output_path):
                self.logger.info("Successfully merged audio with video")
                return True
            else:
                self.logger.error("Failed to merge audio")
                return False
                
        except Exception as e:
            self.logger.error(f"Failed to merge audio: {e}")
            return False
    
    
    def _try_ffmpeg_mp4_from_files(self, frames_list: List[str], output_path: str) -> bool:
        """Try creating MP4 with ffmpeg from frame files"""
        try:
            import subprocess
            import tempfile
            
            if not frames_list:
                self.logger.error("No frames to create video from")
                return False
            
            # Get ffmpeg executable from imageio
            try:
                import imageio_ffmpeg
                ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
                self.logger.debug(f"Using ffmpeg from imageio: {ffmpeg_exe}")
            except Exception as e:
                self.logger.warning(f"Could not find imageio ffmpeg, trying system ffmpeg: {e}")
                ffmpeg_exe = 'ffmpeg'
            
            # Create a temporary file listing frame paths
            self.logger.info(f"Creating frame list file for {len(frames_list)} frames...")
            concat_file = None
            valid_frames = 0
            try:
                with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as f:
                    for i, frame_path in enumerate(frames_list):
                        if os.path.exists(frame_path):
                            abs_path = os.path.abspath(frame_path).replace('\\', '/')
                            f.write(f"file '{abs_path}'\n")
                            valid_frames += 1
                        else:
                            self.logger.warning(f"Frame {i} not found: {frame_path}")
                    concat_file = f.name
            except Exception as e:
                self.logger.error(f"Failed to create frame list: {e}")
                return False
            
            self.logger.info(f"Frame list created with {valid_frames} valid frames")
            
            if not os.path.exists(concat_file):
                self.logger.error("Frame list file was not created")
                return False
            
            self.logger.info(f"Running ffmpeg with frame list: {concat_file}")
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            cmd = [
                ffmpeg_exe,
                '-y',
                '-f', 'concat',
                '-safe', '0',
                '-i', concat_file,
                '-vf', f'scale={self.resolution_width}:{self.resolution_height},fps={self.fps}',
                '-c:v', 'libx264',
                '-pix_fmt', 'yuv420p',
                '-r', str(self.fps),
                '-vsync', 'cfr',
                output_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            # Cleanup concat file
            if os.path.exists(concat_file):
                os.remove(concat_file)
            
            if result.returncode != 0:
                self.logger.error(f"ffmpeg failed: {result.stderr}")
                return False
            
            if not os.path.exists(output_path):
                self.logger.error(f"Output file was not created: {output_path}")
                return False
            
            if os.path.getsize(output_path) <= 1000:
                self.logger.error(f"Output file is too small: {os.path.getsize(output_path)} bytes")
                return False
            
            self.logger.info(f"Successfully created MP4 with ffmpeg: {output_path}")
            return True
            
        except Exception as e:
            self.logger.error(f"ffmpeg MP4 failed: {e}")
            import traceback
            self.logger.debug(traceback.format_exc())
            return False
    
    def _create_video_with_audio_track(self, frames_list: List[str], audio_files: List[str], output_path: str) -> bool:
        """
        Create video with audio track using ffmpeg
        
        Args:
            frames_list: List of frame paths
            audio_files: List of audio paths
            output_path: Output video path
            
        Returns:
            True if successful
        """
        try:
            import subprocess
            import tempfile
            
            # Get ffmpeg executable from imageio
            try:
                import imageio_ffmpeg
                ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
                self.logger.debug(f"Using ffmpeg from imageio: {ffmpeg_exe}")
            except Exception as e:
                self.logger.warning(f"Could not find imageio ffmpeg, trying system ffmpeg: {e}")
                ffmpeg_exe = 'ffmpeg'
            
            self.logger.info(f"Creating temporary video without audio from {len(frames_list)} frames...")
            # First create video without audio
            temp_video = output_path.replace('.mp4', '_temp.mp4')
            
            # Log first few frames for debugging
            if frames_list:
                self.logger.debug(f"First frame: {frames_list[0]}")
                self.logger.debug(f"Last frame: {frames_list[-1] if len(frames_list) > 1 else 'only frame'}")
            
            if not self._try_ffmpeg_mp4_from_files(frames_list, temp_video):
                self.logger.error("Failed to create temp video for audio merging")
                return False
            
            if not os.path.exists(temp_video):
                self.logger.error("Temp video was not created")
                return False
            
            # Log temp video info
            temp_video_size = os.path.getsize(temp_video)
            self.logger.info(f"Temp video created: {temp_video_size} bytes")
            
            # Check temp video with ffprobe if available
            try:
                import imageio_ffmpeg
                ffprobe = imageio_ffmpeg.get_ffmpeg_exe().replace('ffmpeg.exe', 'ffprobe.exe')
                probe_cmd = [
                    ffprobe,
                    '-v', 'error',
                    '-show_entries', 'stream=duration,nb_frames,codec_type',
                    '-of', 'csv=p=0',
                    temp_video
                ]
                probe_result = subprocess.run(probe_cmd, capture_output=True, text=True, timeout=10)
                if probe_result.returncode == 0:
                    self.logger.info(f"Temp video info: {probe_result.stdout.strip()}")
            except Exception as e:
                self.logger.debug(f"Could not probe temp video: {e}")
            
            # Combine audio files into one
            self.logger.info(f"Combining {len(audio_files)} audio files...")
            temp_audio = output_path.replace('.mp4', '_temp_audio.wav')
            audio_list_file = tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8')
            
            for audio_file in audio_files:
                if os.path.exists(audio_file):
                    # Use absolute path and normalize separators
                    abs_path = os.path.abspath(audio_file).replace('\\', '/')
                    audio_list_file.write(f"file '{abs_path}'\n")
            
            audio_list_file.close()
            
            # Concatenate audio files
            self.logger.info(f"Concatenating {len(audio_files)} audio files...")
            
            # Log audio files for debugging
            for i, audio_file in enumerate(audio_files):
                if os.path.exists(audio_file):
                    size = os.path.getsize(audio_file)
                    self.logger.debug(f"Audio {i}: {audio_file} ({size} bytes)")
                else:
                    self.logger.warning(f"Audio {i} not found: {audio_file}")
            
            concat_cmd = [
                ffmpeg_exe, '-y',
                '-f', 'concat',
                '-safe', '0',
                '-i', audio_list_file.name,
                '-acodec', 'pcm_s16le',
                '-ar', '22050',
                temp_audio
            ]
            
            concat_result = subprocess.run(concat_cmd, capture_output=True, text=True, timeout=60)
            
            # Log audio concatenation result
            if concat_result.returncode == 0 and os.path.exists(temp_audio):
                self.logger.info(f"Concatenated audio created: {os.path.getsize(temp_audio)} bytes")
            else:
                self.logger.error(f"Audio concatenation failed: {concat_result.stderr if concat_result.stderr else 'no error message'}")
            
            if concat_result.returncode != 0:
                self.logger.error(f"Audio concatenation failed: {concat_result.stderr}")
                os.remove(audio_list_file.name)
                if os.path.exists(temp_video):
                    os.remove(temp_video)
                return False
            
            if not os.path.exists(temp_audio):
                self.logger.error("Concatenated audio file was not created")
                os.remove(audio_list_file.name)
                if os.path.exists(temp_video):
                    os.remove(temp_video)
                return False
            
            # Merge video and audio
            self.logger.info("Merging video and audio...")
            # Don't use -shortest, let video determine duration
            final_cmd = [
                ffmpeg_exe, '-y',
                '-i', temp_video,
                '-i', temp_audio,
                '-c:v', 'libx264',
                '-c:a', 'aac',
                '-b:a', '192k',
                '-pix_fmt', 'yuv420p',
                '-map', '0:v:0',
                '-map', '1:a:0',
                output_path
            ]
            
            result = subprocess.run(final_cmd, capture_output=True, text=True, timeout=120)
            
            # Log detailed merge results
            if result.stdout:
                self.logger.info(f"Merge command stdout: {result.stdout[:500]}")
            if result.stderr:
                self.logger.info(f"Merge command stderr: {result.stderr[:500]}")
                
            if result.returncode != 0:
                self.logger.error(f"Merge command failed with return code: {result.returncode}")
                if result.stderr:
                    self.logger.error(f"Full stderr: {result.stderr}")
            
            # Check result
            if result.returncode == 0 and os.path.exists(output_path) and os.path.getsize(output_path) > 1000:
                self.logger.info(f"Successfully created MP4 with audio: {output_path}")
                
                # Verify the output has proper duration
                try:
                    import imageio
                    reader = imageio.get_reader(output_path)
                    output_frames = reader.count_frames()
                    self.logger.info(f"Output video has {output_frames} frames")
                    if output_frames < len(frames_list):
                        self.logger.warning(f"Output video has fewer frames ({output_frames}) than expected ({len(frames_list)})")
                except Exception as e:
                    self.logger.debug(f"Could not verify output frames: {e}")
                
                # Cleanup temp files
                if os.path.exists(audio_list_file.name):
                    os.remove(audio_list_file.name)
                if os.path.exists(temp_video):
                    os.remove(temp_video)
                if os.path.exists(temp_audio):
                    os.remove(temp_audio)
                
                return True
            else:
                self.logger.error(f"Audio merge failed: {result.stderr}")
                # Keep temp files for debugging
                self.logger.warning("Keeping temp files for debugging")
            
            return False
            
        except Exception as e:
            self.logger.error(f"Failed to add audio track: {e}")
            import traceback
            self.logger.debug(traceback.format_exc())
            return False
